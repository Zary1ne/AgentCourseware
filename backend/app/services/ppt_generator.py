import os
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from app.config import OUTPUT_DIR
from app.services.llm_service import generate_content


def create_presentation(intent: dict, context: str = "") -> str:
    """根据教学意图生成 PPT 文件，返回文件路径"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 宽屏
    prs.slide_height = Inches(7.5)

    subject = intent.get("subject", "课程")
    topic = intent.get("topic", "教学课件")
    knowledge_points = intent.get("knowledge_points", [])
    objectives = intent.get("objectives", [])
    key_points = intent.get("key_points", [])
    difficult_points = intent.get("difficult_points", [])
    activities = intent.get("activities", [])
    student_profile = intent.get("student_profile")

    # 生成各页内容
    slides_content = _generate_slides_content(intent, context, student_profile)

    # ==== 封面页 ====
    _add_cover_slide(prs, topic, subject)

    # ==== 目录页 ====
    toc_items = ["教学目标", "知识概览", "重点难点"]
    for i, kp in enumerate(knowledge_points[:5]):
        toc_items.append(f"知识点{i+1}: {kp[:30]}")
    toc_items.append("课堂活动")
    toc_items.append("总结与作业")
    _add_toc_slide(prs, toc_items)

    # ==== 教学目标页 ====
    _add_section_slide(prs, "教学目标", objectives or ["掌握核心概念", "理解原理与方法", "能够应用所学知识"])

    # ==== 知识概览页 ====
    _add_section_slide(prs, "知识概览",
                       knowledge_points or ["知识点1", "知识点2", "知识点3"])

    # ==== 内容页（每个知识点一页）====
    for i, kp in enumerate(knowledge_points[:6]):
        detail = slides_content.get(f"kp_{i}", "")
        _add_content_slide(prs, f"知识点{i+1}: {kp}", detail, subject)

    # ==== 重点难点页 ====
    content_parts = []
    if key_points:
        content_parts.append("教学重点：\n" + "\n".join([f"- {p}" for p in key_points]))
    if difficult_points:
        content_parts.append("教学难点：\n" + "\n".join([f"- {p}" for p in difficult_points]))
    _add_section_slide(prs, "重点难点", content_parts)

    # ==== 课堂活动页 ====
    _add_section_slide(prs, "课堂活动",
                       activities or ["小组讨论", "案例分析", "随堂练习"])

    # ==== 互动环节页（如果有测验需求）====
    if intent.get("include_quiz", True):
        _add_quiz_slide(prs, topic, subject, student_profile)

    # ==== 总结页 ====
    summary_items = [
        f"本节课学习了《{topic}》的核心内容",
        f"掌握了 {len(knowledge_points)} 个关键知识点",
    ]
    if objectives:
        summary_items.insert(0, f"教学目标达成：{objectives[0]}")
    _add_summary_slide(prs, summary_items)

    # 保存文件
    filename = f"PPT_{topic}_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)

    return filename


def _generate_slides_content(intent: dict, context: str = "", student_profile: dict = None) -> dict:
    """为每个知识点生成详细内容（优先LLM，失败回退模板）"""
    knowledge_points = intent.get("knowledge_points", [])
    subject = intent.get("subject", "")
    topic = intent.get("topic", "")
    grade = intent.get("grade", "")
    content_map = {}

    for i, kp in enumerate(knowledge_points[:6]):
        try:
            prompt = f"""请为教学课件撰写【{kp}】的详细内容（学科：{subject}，课题：{topic}，年级：{grade}）。
要求：
- 包含核心定义、关键原理、具体示例
- 语言简洁明了，适合PPT展示（每页3-5个要点）
- 每个要点用"• "开头，字数控制在30字以内
- 示例要贴近学生生活实际"""
            llm_content = generate_content(prompt, context, student_profile)
            content_map[f"kp_{i}"] = llm_content
        except Exception:
            # 回退模板
            content_map[f"kp_{i}"] = (
                f"核心概念：{kp}\n\n"
                f"要点一：理解{kp}的基本定义和原理\n"
                f"要点二：掌握{kp}的核心方法和应用场景\n"
                f"要点三：能够运用{kp}解决实际问题\n\n"
                f"示例：结合{subject}学科特点，通过案例加深理解"
            )

    return content_map


# ====== 幻灯片模板函数 ======

def _add_cover_slide(prs, title: str, subtitle: str):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 背景色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x56, 0xDB)  # 蓝色背景
    shape.line.fill.background()

    # 标题
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(1.2))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"{subtitle} | AI智能教学助手"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
    p2.alignment = PP_ALIGN.CENTER

    # 日期
    from datetime import datetime
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10), Inches(0.6))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = datetime.now().strftime("%Y年%m月%d日")
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0x99, 0xBB, 0xEE)
    p3.alignment = PP_ALIGN.CENTER


def _add_toc_slide(prs, items: list):
    """目录页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    _add_slide_title(slide, "目  录")

    # 目录项
    left = Inches(2)
    top = Inches(1.8)
    for i, item in enumerate(items):
        txBox = slide.shapes.add_textbox(left, top + Inches(i * 0.65), Inches(9), Inches(0.55))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"  {i+1:02d}    {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        # 序号左侧色条
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left - Inches(0.1), top + Inches(i * 0.65),
            Inches(0.06), Inches(0.45)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1A, 0x56, 0xDB)
        shape.line.fill.background()


def _add_section_slide(prs, title: str, items: list):
    """章节页（列表型内容）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, title)

    # 分隔线
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.65),
        Inches(2), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x56, 0xDB)
    shape.line.fill.background()

    # 内容项
    top = Inches(2.0)
    for i, item in enumerate(items):
        txBox = slide.shapes.add_textbox(Inches(1.2), top + Inches(i * 0.7), Inches(10.5), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)


def _add_content_slide(prs, title: str, content: str, tag: str = ""):
    """内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, title)

    # 标签
    if tag:
        txBox_tag = slide.shapes.add_textbox(Inches(10.5), Inches(0.35), Inches(2.5), Inches(0.4))
        tf_tag = txBox_tag.text_frame
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = f"[{tag}]"
        p_tag.font.size = Pt(12)
        p_tag.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        p_tag.alignment = PP_ALIGN.RIGHT

    # 分隔线
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.65),
        Inches(2), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x56, 0xDB)
    shape.line.fill.background()

    # 内容
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.5), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = content.strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line.strip()
        if line.strip().startswith("##") or line.strip().startswith("###"):
            p.font.size = Pt(22)
            p.font.bold = True
        elif line.strip().startswith("•") or line.strip().startswith("-"):
            p.font.size = Pt(18)
            p.space_before = Pt(8)
        else:
            p.font.size = Pt(18)
            p.space_before = Pt(6)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def _add_quiz_slide(prs, topic: str, subject: str = "", student_profile: dict = None):
    """互动问答页（优先LLM生成，失败回退模板）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "互动思考")

    try:
        prompt = f"""请为《{topic}》（学科：{subject}）的课堂互动环节设计3个思考题。
要求：
- 问题要有层次：第1题回顾核心概念，第2题联系实际应用，第3题引导深入思考
- 语言贴近学生，开放性而非机械问答
- 每个问题40字以内"""
        llm_output = generate_content(prompt, "", student_profile)
        questions = [q.strip() for q in llm_output.strip().split("\n") if q.strip()][:3]
        if len(questions) < 3:
            raise ValueError("LLM返回问题不足3个")
    except Exception:
        questions = [
            f"回顾一下，{topic}的核心概念是什么？请用自己的话复述。",
            f"你能想到{topic}在生活中的一个应用场景吗？",
            "对本节课的内容，你还有什么疑问？",
        ]

    for i, q in enumerate(questions):
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2 + i * 1.5), Inches(10), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Q{i+1}: {q}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x1A, 0x56, 0xDB)


def _add_summary_slide(prs, items: list):
    """总结页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    shape.line.fill.background()

    _add_slide_title(slide, "课堂总结")

    for i, item in enumerate(items):
        txBox = slide.shapes.add_textbox(Inches(2), Inches(2.2 + i * 0.8), Inches(9), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # 感谢语
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(6.0), Inches(9), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "感谢聆听"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0x1A, 0x56, 0xDB)
    p2.alignment = PP_ALIGN.CENTER


def _add_slide_title(slide, title: str):
    """添加统一的幻灯片标题"""
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(11), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
