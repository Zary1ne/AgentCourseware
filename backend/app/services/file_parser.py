import os
import re
from typing import List, Optional
import PyPDF2
import pdfplumber
from docx import Document
from PIL import Image
from app.config import UPLOAD_DIR


SUPPORTED_TYPES = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "docx",
    ".txt": "text",
    ".md": "text",
    ".jpg": "image",
    ".jpeg": "image",
    ".png": "image",
    ".gif": "image",
    ".webp": "image",
    ".pptx": "pptx",
    ".ppt": "ppt",
}


def get_file_type(filename: str) -> Optional[str]:
    ext = os.path.splitext(filename)[1].lower()
    return SUPPORTED_TYPES.get(ext)


def parse_pdf(filepath: str) -> str:
    text_parts = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"[{i+1}]\n{page_text}")
                tables = page.extract_tables()
                for j, table in enumerate(tables):
                    if table:
                        table_text = "\n".join(
                            [" | ".join([cell or "" for cell in row]) for row in table]
                        )
                        text_parts.append(f"[-{i+1}-{j+1}]\n{table_text}")
    except Exception:
        pass

    if not text_parts:
        try:
            with open(filepath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"[{i+1}]\n{page_text}")
        except Exception as e:
            text_parts.append(f"[PDF: {str(e)}]")

    return "\n\n".join(text_parts)


def parse_docx(filepath: str) -> str:
    text_parts = []
    try:
        doc = Document(filepath)
        for para in doc.paragraphs:
            if para.text.strip():
                style = para.style.name if para.style else "Normal"
                if "Heading" in style:
                    text_parts.append(f"## {para.text}")
                else:
                    text_parts.append(para.text)
        for i, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text for cell in row.cells]
                rows.append(" | ".join(cells))
            text_parts.append(f"\n[{i+1}]\n" + "\n".join(rows))
    except Exception as e:
        text_parts.append(f"[Word: {str(e)}]")

    return "\n\n".join(text_parts)


def parse_txt(filepath: str) -> str:
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(filepath, "r", encoding="gbk") as f:
            return f.read()


def parse_image(filepath: str) -> str:
    try:
        img = Image.open(filepath)
        info = {
            "": img.format,
            "": f"{img.width}x{img.height}",
            "": img.mode,
        }
        description = f"[{os.path.basename(filepath)}]\n"
        description += f": {info['']}, : {info['']}\n"
        description += "[]"
        return description
    except Exception as e:
        return f"[: {str(e)}]"


def parse_pptx(filepath: str) -> str:
    text_parts = []
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text)
            if slide_texts:
                text_parts.append(f"[{i+1}]\n" + "\n".join(slide_texts))
    except Exception as e:
        text_parts.append(f"[PPT: {str(e)}]")

    return "\n\n".join(text_parts)


def parse_ppt(filepath: str) -> str:
    text_parts = []
    try:
        with open(filepath, "rb") as f:
            data = f.read()
        decoded = data.decode("utf-16-le", errors="ignore")
        fragments = re.findall(r"[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef\u0020-\u007e]{3,}", decoded)
        seen = set()
        for frag in fragments:
            frag = frag.strip()
            if not frag or frag in seen:
                continue
            if re.fullmatch(r"[\s\d\W]+", frag):
                continue
            seen.add(frag)
            text_parts.append(frag)
        if not text_parts:
            return "[] .pptx "
        return "[]\n" + "\n".join(text_parts)
    except Exception as e:
        return f"[PPT: {str(e)}]"


def parse_file(filepath: str, filename: str) -> dict:
    file_type = get_file_type(filename)

    parsers = {
        "pdf": parse_pdf,
        "docx": parse_docx,
        "text": parse_txt,
        "image": parse_image,
        "pptx": parse_pptx,
        "ppt": parse_ppt,
    }

    parser = parsers.get(file_type)
    if not parser:
        return {"success": False, "content": f": {file_type}", "file_type": file_type}

    try:
        content = parser(filepath)
        if len(content) > 10000:
            content = content[:10000] + "\n\n[...]"
        return {"success": True, "content": content, "file_type": file_type}
    except Exception as e:
        return {"success": False, "content": f": {str(e)}", "file_type": file_type}